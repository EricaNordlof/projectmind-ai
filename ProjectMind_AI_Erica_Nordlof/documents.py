from __future__ import annotations

import hashlib
import html
import io
import json
import re
from dataclasses import dataclass
from datetime import datetime, timezone
from html.parser import HTMLParser
from pathlib import Path
from typing import Any

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from config import settings
from db import connection, execute, fetchall, fetchone, sql
from storage import read_bytes


TEXT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".html", ".htm", ".xml", ".yaml", ".yml",
    ".toml", ".ini", ".cfg", ".py", ".java", ".js", ".mjs", ".cjs", ".ts",
    ".tsx", ".jsx", ".css", ".scss", ".sql", ".sh", ".bat", ".ps1", ".env",
    ".log", ".rst",
}
SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | {".pdf", ".docx", ".xlsx", ".pptx"}

SOURCE_PATTERN = re.compile(
    r"<<< KÄLLA: (?P<filename>.*?) \| (?P<locator>.*?) >>>\n(?P<text>.*?)\n<<< SLUT KÄLLA: .*? >>>",
    re.DOTALL,
)

SWEDISH_STOPWORDS = {
    "och", "att", "det", "den", "detta", "de", "dem", "en", "ett", "är", "var", "vad",
    "som", "för", "på", "i", "av", "till", "med", "om", "från", "har", "kan", "ska", "skall",
    "jag", "du", "vi", "ni", "min", "mitt", "mina", "din", "ditt", "dina", "alla", "hela",
    "filer", "fil", "dokument", "granska", "analysera", "sammanfatta", "visa", "säg", "hur",
}


@dataclass
class ExtractionResult:
    text: str
    locator_count: int
    error: str = ""


class VisibleTextParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []
        self.hidden_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        tag = tag.lower()
        if tag in {"script", "style", "noscript", "svg"}:
            self.hidden_depth += 1
        elif tag in {"p", "div", "br", "li", "tr", "h1", "h2", "h3", "h4", "h5", "h6"}:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        tag = tag.lower()
        if tag in {"script", "style", "noscript", "svg"} and self.hidden_depth:
            self.hidden_depth -= 1
        elif tag in {"p", "div", "li", "tr", "h1", "h2", "h3", "h4", "h5", "h6"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self.hidden_depth:
            self.parts.append(data)


def clean_text(value: str) -> str:
    value = value.replace("\x00", " ")
    value = value.replace("\r\n", "\n").replace("\r", "\n")
    value = re.sub(r"[ \t]+", " ", value)
    value = re.sub(r"\n[ \t]+", "\n", value)
    value = re.sub(r"\n{4,}", "\n\n\n", value)
    return value.strip()


def decode_text(raw: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def source_block(filename: str, locator: str, text: str) -> str:
    return (
        f"<<< KÄLLA: {filename} | {locator} >>>\n"
        f"{clean_text(text)}\n"
        f"<<< SLUT KÄLLA: {filename} | {locator} >>>"
    )


def extract_pdf(raw: bytes, filename: str) -> ExtractionResult:
    reader = PdfReader(io.BytesIO(raw))
    blocks: list[str] = []
    for page_number, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text(extraction_mode="layout") or page.extract_text() or ""
        except TypeError:
            text = page.extract_text() or ""
        cleaned = clean_text(text)
        if not cleaned:
            cleaned = "[Ingen maskinläsbar text hittades. Sidan kan vara skannad eller bildbaserad.]"
        blocks.append(source_block(filename, f"sida {page_number}", cleaned))
    return ExtractionResult("\n\n".join(blocks), len(reader.pages))


def extract_docx(raw: bytes, filename: str) -> ExtractionResult:
    document = Document(io.BytesIO(raw))
    parts: list[str] = []
    for paragraph in document.paragraphs:
        text = clean_text(paragraph.text)
        if text:
            parts.append(text)
    for table_index, table in enumerate(document.tables, start=1):
        rows: list[str] = []
        for row in table.rows:
            cells = [clean_text(cell.text) for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"[Tabell {table_index}]\n" + "\n".join(rows))
    return ExtractionResult(source_block(filename, "dokument", "\n\n".join(parts)), 1)


def extract_xlsx(raw: bytes, filename: str) -> ExtractionResult:
    workbook = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    blocks: list[str] = []
    for worksheet in workbook.worksheets:
        rows: list[str] = []
        char_count = 0
        for row in worksheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            while values and not values[-1]:
                values.pop()
            if any(values):
                line = " | ".join(values)
                rows.append(line)
                char_count += len(line)
            if char_count >= settings.smart_max_file_chars:
                rows.append("[Bladet trunkerades eftersom det är mycket stort.]")
                break
        blocks.append(source_block(filename, f"blad {worksheet.title}", "\n".join(rows)))
    return ExtractionResult("\n\n".join(blocks), len(workbook.worksheets))


def extract_pptx(raw: bytes, filename: str) -> ExtractionResult:
    presentation = Presentation(io.BytesIO(raw))
    blocks: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = clean_text(getattr(shape, "text", ""))
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                rows = [
                    " | ".join(clean_text(cell.text) for cell in row.cells)
                    for row in shape.table.rows
                ]
                if rows:
                    parts.append("\n".join(rows))
        blocks.append(source_block(filename, f"bild {slide_number}", "\n\n".join(parts)))
    return ExtractionResult("\n\n".join(blocks), len(presentation.slides))


def extract_html(raw: bytes, filename: str) -> ExtractionResult:
    parser = VisibleTextParser()
    parser.feed(decode_text(raw))
    text = html.unescape("".join(parser.parts))
    return ExtractionResult(source_block(filename, "dokument", text), 1)


def extract_json(raw: bytes, filename: str) -> ExtractionResult:
    text = decode_text(raw)
    try:
        text = json.dumps(json.loads(text), ensure_ascii=False, indent=2)
    except json.JSONDecodeError:
        pass
    return ExtractionResult(source_block(filename, "dokument", text), 1)


def extract_document(raw: bytes, filename: str, mime_type: str = "") -> ExtractionResult:
    extension = Path(filename).suffix.lower()
    if extension == ".pdf" or mime_type == "application/pdf":
        return extract_pdf(raw, filename)
    if extension == ".docx":
        return extract_docx(raw, filename)
    if extension == ".xlsx":
        return extract_xlsx(raw, filename)
    if extension == ".pptx":
        return extract_pptx(raw, filename)
    if extension in {".html", ".htm"}:
        return extract_html(raw, filename)
    if extension == ".json":
        return extract_json(raw, filename)
    if extension in TEXT_EXTENSIONS or mime_type.startswith("text/"):
        return ExtractionResult(source_block(filename, "dokument", decode_text(raw)), 1)
    return ExtractionResult("", 0, "Filtypen kan lagras men har ingen textutvinning.")


def index_file(record: dict[str, Any], force: bool = False) -> dict[str, Any]:
    try:
        raw = read_bytes("uploads", record["project_id"], record["stored_name"])
    except Exception as exc:
        return {"text": "", "locator_count": 0, "error": f"Filen kunde inte läsas: {exc}"}

    digest = hashlib.sha256(raw).hexdigest()
    cached = fetchone("SELECT * FROM document_text_cache WHERE file_id=?", (record["id"],))
    if not force and cached and cached.get("content_sha256") == digest:
        return {
            "text": cached.get("extracted_text", ""),
            "locator_count": int(cached.get("locator_count", 0) or 0),
            "error": cached.get("extraction_error", ""),
            "extracted_at": cached.get("extracted_at", ""),
        }

    try:
        result = extract_document(raw, record["original_name"], record.get("mime_type", "") or "")
        text = result.text[: settings.smart_max_file_chars]
        error = result.error
        if len(result.text) > settings.smart_max_file_chars:
            error = (error + " " if error else "") + "Texten trunkerades vid filgränsen."
        locator_count = result.locator_count
    except Exception as exc:
        text = ""
        locator_count = 0
        error = f"Textutvinningen misslyckades: {type(exc).__name__}: {exc}"

    extracted_at = datetime.now(timezone.utc).isoformat(timespec="seconds")
    with connection() as conn:
        conn.execute(sql("DELETE FROM document_text_cache WHERE file_id=?"), (record["id"],))
        conn.execute(
            sql(
                """
                INSERT INTO document_text_cache(
                    file_id,content_sha256,extracted_text,locator_count,extraction_error,extracted_at
                ) VALUES(?,?,?,?,?,?)
                """
            ),
            (record["id"], digest, text, locator_count, error, extracted_at),
        )

    return {
        "text": text,
        "locator_count": locator_count,
        "error": error,
        "extracted_at": extracted_at,
    }


def reindex_project(project_id: str, force: bool = True) -> list[dict[str, Any]]:
    results: list[dict[str, Any]] = []
    for record in fetchall(
        "SELECT * FROM project_files WHERE project_id=? ORDER BY created_at DESC",
        (project_id,),
    ):
        result = index_file(record, force=force)
        results.append({"file": record["original_name"], **result})
    return results


def file_index_status(file_id: str) -> dict[str, Any] | None:
    return fetchone("SELECT * FROM document_text_cache WHERE file_id=?", (file_id,))


def revision_priority(filename: str, created_at: str = "") -> tuple[int, str]:
    name = filename.casefold()
    score = 0
    if any(token in name for token in ("revider", " revision", " rev ", "-rev", "_rev")):
        score += 100
    for match in re.findall(r"(?:utgåva|version|ver|rev)[ _.-]*(\d+)", name):
        score += min(int(match), 20) * 5
    if re.search(r"\b(20\d{4}|\d{6})\b", name):
        score += 15
    if "kursplan" in name or name.startswith("kp"):
        score += 10
    if "schema" in name:
        score += 8
    if "välkomst" in name:
        score += 6
    return score, created_at or ""


def query_tokens(query: str) -> set[str]:
    words = {
        word.casefold()
        for word in re.findall(r"[A-Za-zÅÄÖåäö0-9_-]{3,}", query)
    }
    return {word for word in words if word not in SWEDISH_STOPWORDS}


def broad_request(query: str) -> bool:
    normalized = query.casefold()
    phrases = (
        "alla filer", "alla dokument", "granska allt", "granska dem", "samlad bedömning",
        "helhetsbedömning", "hela projektet", "sammanfatta allt", "jämför dokumenten",
        "gå igenom alla", "granska dom", "granska de", "vad innehåller projektet",
    )
    return any(phrase in normalized for phrase in phrases)


def split_blocks(extracted_text: str) -> list[dict[str, str]]:
    blocks: list[dict[str, str]] = []
    for match in SOURCE_PATTERN.finditer(extracted_text):
        blocks.append(
            {
                "filename": match.group("filename"),
                "locator": match.group("locator"),
                "text": match.group("text").strip(),
                "full": match.group(0),
            }
        )
    if not blocks and extracted_text.strip():
        blocks.append({"filename": "okänd", "locator": "dokument", "text": extracted_text, "full": extracted_text})
    return blocks


def score_block(block: dict[str, str], tokens: set[str], filename: str, priority: int) -> int:
    if not tokens:
        return priority
    haystack = (filename + " " + block["locator"] + " " + block["text"]).casefold()
    score = priority // 10
    for token in tokens:
        count = haystack.count(token)
        if count:
            score += 8 + min(count, 8) * 2
        if token in filename.casefold():
            score += 12
    return score


def build_project_context(project_id: str | None, query: str = "") -> str:
    if not project_id:
        return ""
    project = fetchone("SELECT * FROM projects WHERE id=?", (project_id,))
    if not project:
        return ""

    files = fetchall(
        """
        SELECT f.*,c.extracted_text,c.locator_count,c.extraction_error,c.extracted_at
        FROM project_files f
        LEFT JOIN document_text_cache c ON c.file_id=f.id
        WHERE f.project_id=?
        ORDER BY f.created_at DESC
        """,
        (project_id,),
    )

    enriched: list[dict[str, Any]] = []
    for record in files:
        if record.get("extracted_text") is None:
            indexed = index_file(record)
            record.update(
                extracted_text=indexed.get("text", ""),
                locator_count=indexed.get("locator_count", 0),
                extraction_error=indexed.get("error", ""),
                extracted_at=indexed.get("extracted_at", ""),
            )
        enriched.append(record)

    enriched.sort(
        key=lambda item: revision_priority(item["original_name"], item.get("created_at", "")),
        reverse=True,
    )

    versions = fetchall(
        "SELECT version_label,comment,original_name,created_at FROM versions WHERE project_id=? ORDER BY created_at DESC LIMIT 100",
        (project_id,),
    )

    header = [
        "PROJEKTMETADATA",
        f"Projekt: {project['name']}",
        f"Status: {project['status']}",
        f"Stack: {project['stack']}",
        f"Beskrivning: {project['description']}",
        f"Anteckningar: {project['notes']}",
        "",
        "DOKUMENTKATALOG",
    ]
    for item in enriched:
        status = "text indexerad"
        if item.get("extraction_error"):
            status = item["extraction_error"]
        locator_count = int(item.get("locator_count", 0) or 0)
        locator_info = f", {locator_count} källavsnitt" if locator_count else ""
        header.append(
            f"- {item['original_name']} ({item.get('size_bytes', 0)} bytes{locator_info}) — {status}"
        )

    header.extend(["", "VERSIONER/ZIP-ARKIV"])
    for version in versions:
        header.append(
            f"- {version['version_label']}: {version['original_name']} — {version['comment']} ({version['created_at']})"
        )
    header.extend(
        [
            "",
            "KÄLLMATERIAL",
            "Reviderade filer och högre utgåvor visas normalt före äldre filer. Det är en prioriteringssignal, inte automatiskt ett bevis på giltighet.",
        ]
    )

    context = "\n".join(header)
    remaining = settings.smart_context_max_chars - len(context)
    tokens = query_tokens(query)
    include_all = broad_request(query) or sum(len(item.get("extracted_text") or "") for item in enriched) <= remaining

    candidates: list[tuple[int, int, str, str]] = []
    for file_order, item in enumerate(enriched):
        extracted_text = item.get("extracted_text") or ""
        priority, _ = revision_priority(item["original_name"], item.get("created_at", ""))
        for block_order, block in enumerate(split_blocks(extracted_text)):
            score = 100000 - file_order * 100 - block_order if include_all else score_block(
                block, tokens, item["original_name"], priority
            )
            candidates.append((score, -file_order, item["original_name"], block["full"]))

    candidates.sort(key=lambda item: (item[0], item[1]), reverse=True)
    included_blocks = 0
    included_files: set[str] = set()
    omitted_files: set[str] = set()

    for score, _, filename, block in candidates:
        if not include_all and score <= 0 and included_blocks >= 8:
            omitted_files.add(filename)
            continue
        addition = f"\n\n===== DOKUMENT: {filename} =====\n{block}"
        if len(addition) <= remaining:
            context += addition
            remaining -= len(addition)
            included_blocks += 1
            included_files.add(filename)
        else:
            omitted_files.add(filename)
        if remaining < 2500:
            break

    for item in enriched:
        if item.get("extracted_text") and item["original_name"] not in included_files:
            omitted_files.add(item["original_name"])

    context += f"\n\nKONTEXTSTATUS: {included_blocks} källavsnitt från {len(included_files)} dokument inkluderades."
    if omitted_files:
        context += " Material som inte rymdes eller inte bedömdes mest relevant: " + ", ".join(sorted(omitted_files))
    return context
